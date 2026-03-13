import os
import asyncio
import gc
import json
import shutil
import pandas as pd
from typing import List, Dict, Any
from pathlib import Path

from docx import Document as DocxDocument        
from pptx import Presentation        
import pytesseract
from PIL import Image
from pdf2image import convert_from_path

from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_groq import ChatGroq
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_core.documents import Document

CHROMA_DIR = "/tmp/chroma_db"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".txt",
    ".xlsx", ".xls", ".csv",
    ".png", ".jpg", ".jpeg", ".tiff", ".bmp"
}

# ── Improvement 3: Better chunking separators ──────────────────────────────────
CHUNK_SIZE    = 1000
CHUNK_OVERLAP = 150
SEPARATORS    = ["\n\n", "\n", ". ", "! ", "? ", ", ", " ", ""]

SYSTEM_PROMPT = """You are an expert document analyst.
Answer questions based ONLY on the provided context below.
If the answer is not found in the context, say exactly:
"I couldn't find this information in the uploaded documents."

Be precise, structured, and always reference page numbers when available.

Context:
{context}
"""

SUMMARY_PROMPT = """You are a document analyst. Given the beginning of a document, produce:
1. A 2-3 sentence summary of what this document is about
2. 3-5 key topics covered (as a comma-separated list)
3. 3 suggested questions a user might ask about this document

Respond in this exact JSON format:
{{
  "summary": "...",
  "topics": ["topic1", "topic2", "topic3"],
  "suggested_questions": ["question1", "question2", "question3"]
}}

Document content:
{content}
"""

class RAGPipeline:
    def __init__(self):
        print("Initializing RAG Pipeline...")

        self.embeddings = GoogleGenerativeAIEmbeddings(
            model="models/gemini-embedding-001",
            google_api_key=os.environ.get("GOOGLE_API_KEY")
        )

        groq_key = os.environ.get("GROQ_API_KEY")
        if not groq_key:
            raise ValueError("GROQ_API_KEY not found. Check your .env file.")

        self.llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0.1,
            api_key=groq_key,
        )

        self.vectorstore = None
        self.indexed_docs: List[str] = []
        
        # ── Improvement 1: Conversation memory per session ────────────────────
        self.conversation_history: List[Dict[str, str]] = []
        
        self._load_existing_index()
        print("RAG Pipeline ready.")

    def _load_existing_index(self):
        """Load ChromaDB index if exists."""

        try:
            self.vectorstore = Chroma(
                persist_directory=CHROMA_DIR,
                embedding_function=self.embeddings
            )

            count = self.vectorstore._collection.count()

            if count == 0:
                self.vectorstore = None
            else:
                print(f"Loaded existing index ({count} vectors).")

        except Exception as e:
            print("Could not load index:", e)
            self.vectorstore = None
            
    # ── Sanitize filename ─────────────────────────────────────────────────────
    def _sanitize_filename(self, file_path: str) -> str:
        original_path = Path(file_path)
        raw_name  = original_path.name
        safe_name = "".join(
            c if (c.isalnum() or c in "._- ") else "_"
            for c in raw_name
        )
        safe_name = safe_name.strip("_ ") or "document"
        safe_path = original_path.parent / safe_name
        if safe_path != original_path:
            shutil.copy2(str(original_path), str(safe_path))
            print(f"   Renamed to safe path: {safe_name}")
            return str(safe_path)
        return file_path
    
    # ── Universal index method ────────────────────────────────────────────────
    def index_file(self, file_path: str) -> Dict[str, Any]:
        """
        Universal loader — detects file type and indexes accordingly.
        Supports: PDF, DOCX, PPTX, TXT, XLSX, XLS, CSV, Images (OCR)
        """
        file_path = self._sanitize_filename(file_path)
        ext       = Path(file_path).suffix.lower()

        print(f"Loading file: {file_path} (type: {ext})")

        if ext == ".pdf":
            pages = self._load_pdf(file_path)
        elif ext == ".docx":
            pages = self._load_docx(file_path)
        elif ext == ".pptx":
            pages = self._load_pptx(file_path)
        elif ext == ".txt":
            pages = self._load_txt(file_path)
        elif ext in (".xlsx", ".xls"):
            pages = self._load_excel(file_path)
        elif ext == ".csv":
            pages = self._load_csv(file_path)
        elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
            pages = self._load_image_ocr(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

        pages = [p for p in pages if p.page_content and p.page_content.strip()]
        print(f"   {len(pages)} pages/sections loaded")

        if not pages:
            raise ValueError("No text could be extracted from this file.")

        return self._chunk_and_index(file_path, pages)

    # ── Backward compatible ───────────────────────────────────────────────────
    def index_pdf(self, file_path: str) -> Dict[str, Any]:
        """Backward compatible — calls universal index_file."""
        return self.index_file(file_path)
    
    # ── PDF Loader ────────────────────────────────────────────────────────────
    def _load_pdf(self, file_path: str) -> List[Document]:
        loader    = PyPDFLoader(file_path)
        pages     = loader.load()
        non_empty = [p for p in pages if p.page_content.strip()]
        if not non_empty:
            print("   Scanned PDF detected — applying OCR...")
            return self._ocr_pdf(file_path)
        return pages
    
    # ── DOCX Loader ───────────────────────────────────────────────────────────
    def _load_docx(self, file_path: str) -> List[Document]:
        doc  = DocxDocument(file_path)
        text = "\n\n".join([
            para.text for para in doc.paragraphs
            if para.text.strip()
        ])
        return [Document(
            page_content=text,
            metadata={"source_file": Path(file_path).name, "page": 1}
        )]

    # ── PPTX Loader ───────────────────────────────────────────────────────────
    def _load_pptx(self, file_path: str) -> List[Document]:
        prs   = Presentation(file_path)
        pages = []
        for i, slide in enumerate(prs.slides):
            text = "\n".join([
                shape.text for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ])
            if text.strip():
                pages.append(Document(
                    page_content=text,
                    metadata={"source_file": Path(file_path).name, "page": i + 1}
                ))
        return pages

    # ── TXT Loader ────────────────────────────────────────────────────────────
    def _load_txt(self, file_path: str) -> List[Document]:
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
        return [Document(
            page_content=text,
            metadata={"source_file": Path(file_path).name, "page": 1}
        )]

    # ── Excel Loader ──────────────────────────────────────────────────────────
    def _load_excel(self, file_path: str) -> List[Document]:
        xl    = pd.ExcelFile(file_path)
        pages = []
        for sheet_name in xl.sheet_names:
            df   = xl.parse(sheet_name)
            text = f"Sheet: {sheet_name}\n\n"
            text += f"Columns: {', '.join(df.columns.astype(str).tolist())}\n\n"
            text += df.to_string(index=False)
            pages.append(Document(
                page_content=text,
                metadata={
                    "source_file": Path(file_path).name,
                    "sheet":       sheet_name,
                    "page":        1,
                    "rows":        len(df),
                    "columns":     len(df.columns)
                }
            ))
        return pages

    # ── CSV Loader ────────────────────────────────────────────────────────────
    def _load_csv(self, file_path: str) -> List[Document]:
        df   = pd.read_csv(file_path)
        text = f"Columns: {', '.join(df.columns.astype(str).tolist())}\n\n"
        text += df.to_string(index=False)
        return [Document(
            page_content=text,
            metadata={
                "source_file": Path(file_path).name,
                "page":    1,
                "rows":    len(df),
                "columns": len(df.columns)
            }
        )]

    # ── Image OCR Loader ──────────────────────────────────────────────────────
    def _load_image_ocr(self, file_path: str) -> List[Document]:
        try:
            img  = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return [Document(
                page_content=text,
                metadata={"source_file": Path(file_path).name, "page": 1}
            )]
        except Exception as e:
            raise ValueError(f"OCR failed: {e}. Make sure tesseract is installed.")

    # ── OCR for scanned PDFs ──────────────────────────────────────────────────
    def _ocr_pdf(self, file_path: str) -> List[Document]:
        try:
            images = convert_from_path(file_path)
            pages  = []
            for i, img in enumerate(images):
                text = pytesseract.image_to_string(img)
                if text.strip():
                    pages.append(Document(
                        page_content=text,
                        metadata={"source_file": Path(file_path).name, "page": i + 1}
                    ))
            return pages
        except Exception as e:
            raise ValueError(
                f"OCR failed: {e}. "
                f"Run: apt install tesseract-ocr poppler-utils"
            )
            
    # ── Chunk + Index (shared for all formats) ────────────────────────────────
    def _chunk_and_index(self, file_path: str, pages: List[Document]) -> Dict[str, Any]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            separators=SEPARATORS,
            length_function=len,
        )
        chunks = splitter.split_documents(pages)
        chunks = [c for c in chunks if c.page_content and c.page_content.strip()]
        print(f"   {len(chunks)} chunks created")

        if not chunks:
            raise ValueError("No content could be extracted from this file.")

        filename = Path(file_path).name
        for i, chunk in enumerate(chunks):
            chunk.metadata["source_file"] = filename
            chunk.metadata["chunk_index"] = i
            chunk.metadata["total_chunks"] = len(chunks)
            chunk.metadata["position"]     = f"chunk {i+1}/{len(chunks)}"

        if self.vectorstore is None:
            self.vectorstore = Chroma.from_documents(
                documents=chunks,
                embedding=self.embeddings,
                persist_directory=CHROMA_DIR
            )
        else:
            self.vectorstore.add_documents(chunks)

        first_content = " ".join([p.page_content for p in pages[:3]])[:3000]
        doc_info = self._generate_summary(filename, first_content, len(chunks))

        self.indexed_docs = [d for d in self.indexed_docs if d["name"] != filename]
        self.indexed_docs.append(doc_info)

        return doc_info

    def _generate_summary(self, filename: str, content: str, chunks: int) -> Dict:
        """Generate auto summary, topics and suggested questions."""
        import json
        try:
            prompt = ChatPromptTemplate.from_messages([
                ("human", SUMMARY_PROMPT)
            ])
            chain = prompt | self.llm | StrOutputParser()
            raw = chain.invoke({"content": content})

            # Strip markdown fences if present
            raw = raw.strip()
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            parsed = json.loads(raw.strip())

            return {
                "name": filename,
                "chunks": chunks,
                "summary": parsed.get("summary", ""),
                "topics": parsed.get("topics", []),
                "suggested_questions": parsed.get("suggested_questions", []),
            }
        except Exception as e:
            print(f"Summary generation failed: {e}")
            return {
                "name": filename,
                "chunks": chunks,
                "summary": "",
                "topics": [],
                "suggested_questions": [],
            }
        
    # ── Improvement 1: Multi-turn conversation memory ─────────────────────────
    async def answer(self, question: str) -> Dict[str, Any]:
        if self.vectorstore is None:
            return {
                "answer": "No documents indexed yet. Please upload a PDF first.",
                "sources": [], "confidence": 0, "conversation_turn": 0
            }

        loop = asyncio.get_event_loop()

        retriever = self.vectorstore.as_retriever(
            search_type="mmr",
            search_kwargs={"k": 5, "fetch_k": 10}
        )

        relevant_docs = await loop.run_in_executor(
            None, retriever.invoke, question
        )
        print(f"{len(relevant_docs)} relevant chunks retrieved")

        if not relevant_docs:
            return {
                "answer": "I couldn't find relevant information in the uploaded documents.",
                "sources": [], "confidence": 0, 
                "conversation_turn": len(self.conversation_history)
            }

        # Improvement 5: Include paragraph/position in context
        context = "\n\n---\n\n".join([
            f"[Source: {doc.metadata.get('source_file','unknown')} | "
            f"Page {doc.metadata.get('page', 0) + 1} | "
            f"Position: {doc.metadata.get('position', 'unknown')}]\n{doc.page_content}"
            for doc in relevant_docs
        ])

        # Improvement 1: Build messages with conversation history
        messages = [("system", SYSTEM_PROMPT)]

        # Add last 4 turns of history (keep context window manageable)
        for turn in self.conversation_history[-4:]:
            messages.append(("human", turn["question"]))
            messages.append(("assistant", turn["answer"]))

        messages.append(("human", "{question}"))

        prompt = ChatPromptTemplate.from_messages(messages)
        chain = prompt | self.llm | StrOutputParser()

        answer_text = await loop.run_in_executor(
            None,
            lambda: chain.invoke({"context": context, "question": question})
        )

        # Save to conversation history
        self.conversation_history.append({
            "question": question,
            "answer": answer_text
        })

        # Improvement 5: Rich source info with page + position
        sources = list({
            f"{doc.metadata.get('source_file','unknown')} — Page {doc.metadata.get('page', 0) + 1}"
            for doc in relevant_docs
        })

        confidence = min(len(relevant_docs) * 20, 95)

        return {
            "answer": answer_text,
            "sources": sources,
            "confidence": confidence,
            "conversation_turn": len(self.conversation_history),
            "source_passages": [
                {
                    "file": doc.metadata.get("source_file", "unknown"),
                    "page": doc.metadata.get("page", 0) + 1,
                    "position": doc.metadata.get("position", ""),
                    "text": doc.page_content[:200] + "..."
                }
                for doc in relevant_docs[:3]
            ]
        }    
    
    def answer_sync(self, question: str) -> dict:
        """Fully synchronous version — safe to run in executor."""

        print(f"DEBUG: answer_sync called: {question}")

        if self.vectorstore is None:
            return {
                "answer": "No documents indexed yet. Please upload a PDF first.",
                "sources": [],
                "confidence": 0
            }

        retriever = self.vectorstore.as_retriever(
            search_type="mmr",
            search_kwargs={"k": 5, "fetch_k": 10}
        )

        relevant_docs = retriever.invoke(question)
        print(f"DEBUG: retrieved {len(relevant_docs)} docs")

        if not relevant_docs:
            return {
                "answer": "I couldn't find relevant information in the uploaded documents.",
                "sources": [],
                "confidence": 0
            }

        context = "\n\n---\n\n".join([
            f"[Source: {doc.metadata.get('source_file', 'unknown')} | "
            f"Page {doc.metadata.get('page', 0) + 1}]\n{doc.page_content}"
            for doc in relevant_docs
        ])
        prompt = ChatPromptTemplate.from_messages([
            ("system", SYSTEM_PROMPT),
            ("human", "{question}")
        ])
        chain = prompt | self.llm | StrOutputParser()
        answer_text = chain.invoke({"context": context, "question": question})

        print(f"DEBUG: got answer: {answer_text[:80]}...")

        sources = list({
            f"{doc.metadata.get('source_file', 'unknown')} — Page {doc.metadata.get('page', 0) + 1}"
            for doc in relevant_docs
        })
        confidence = min(len(relevant_docs) * 20, 95)

        return {
            "answer": answer_text,
            "sources": sources,
            "confidence": confidence
        }
        
    def clear_conversation(self):
        """Clear conversation history only — keep documents indexed."""
        self.conversation_history = []

    def list_indexed_documents(self) -> List[str]:
        return self.indexed_docs

    def clear_index(self):
        """Clear vectorstore safely without deleting filesystem."""

        print("Clearing index...")

        try:
            if self.vectorstore:
                self.vectorstore.delete_collection()
        except Exception as e:
            print("Collection delete warning:", e)

        self.vectorstore = None
        self.indexed_docs = []
        self.conversation_history = []
        gc.collect()

        print("Index cleared.")
